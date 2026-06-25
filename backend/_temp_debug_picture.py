from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
prs = Presentation('temp_test_output.pptx')
slide = prs.slides[0]
for idx, shape in enumerate(slide.shapes):
    ph = getattr(shape, 'placeholder_format', None)
    ph_type = getattr(ph, 'type', None)
    print('shape', idx, 'name:', shape.name, 'shape_type:', shape.shape_type, 'ph_type:', ph_type)
    if ph is not None:
        print('  placeholder attrs', [a for a in dir(ph) if 'insert' in a.lower() or 'picture' in a.lower()])
    if hasattr(shape, 'image'):
        try:
            img = shape.image
            print('  has image, size', img.size)
        except Exception as e:
            print('  has image attr but error', e)
    if hasattr(shape, 'has_image'):
        print('  has_image', shape.has_image)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print('  picture shape found')
