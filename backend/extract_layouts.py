import json
import os
import sys
from pptx import Presentation

def extract_layouts(directory):
    print(f"Scanning directory: {directory}")
    for file in os.listdir(directory):
        if file.endswith(".potx") or file.endswith(".pptx"):
            path = os.path.join(directory, file)
            print(f"Processing: {file}")
            try:
                prs = Presentation(path)
            except Exception as e:
                print(f"Failed to load {file}: {e}")
                continue
            
            layouts = []
            for i, layout in enumerate(prs.slide_layouts):
                placeholders = []
                for ph in layout.placeholders:
                    ptype = getattr(ph.placeholder_format, "type", None)
                    placeholders.append({
                        "name": ph.name,
                        "type": ptype
                    })
                layouts.append({
                    "layout_index": i,
                    "name": layout.name,
                    "placeholders": placeholders
                })
                
            out_name = os.path.splitext(file)[0] + ".json"
            out_path = os.path.join(directory, out_name)
            with open(out_path, "w", encoding="utf-8") as f:
                json.dump({"layouts": layouts}, f, indent=2)
            print(f"Saved: {out_name}")

if __name__ == "__main__":
    target_dir = os.path.join(os.path.dirname(__file__), "sample_ppt")
    extract_layouts(target_dir)
