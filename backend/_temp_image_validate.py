from backend.components.create_ppt import create_ppt
from pptx import Presentation
slides = [
    {"title": "Cover Slide", "content": ["Welcome"], "content_type": "title", "image_index": None},
    {"title": "Image Slide", "content": ["This is a test slide with image."], "content_type": "image_text", "image_index": 0}
]
audio_folder = 'backend'
output_ppt = 'backend/temp_test_output.pptx'
images_folder = 'backend/projects/Taxation/v1/slides_images'
count = create_ppt(slides, audio_folder, output_ppt=output_ppt, template_name='template2.potx', images_folder=images_folder)
print('created', count, output_ppt)
prs = Presentation(output_ppt)
for i, slide in enumerate(prs.slides, start=1):
    print('Slide', i, 'shapes', len(slide.shapes))
    for shape in slide.shapes:
        ph = getattr(shape, 'placeholder_format', None)
        ph_type = getattr(ph, 'type', None)
        print('  name=', shape.name, 'shape_type=', shape.shape_type, 'ph_type=', ph_type)
        if shape.shape_type.name == 'PICTURE':
            print('    picture shape present')
        if hasattr(shape, 'image'):
            try:
                print('    image size', shape.image.size)
            except Exception:
                pass
