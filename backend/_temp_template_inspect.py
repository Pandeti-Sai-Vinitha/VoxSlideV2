from pathlib import Path
from pptx import Presentation
root = Path('sample_ppt')
for template in sorted(root.glob('*.potx')):
    prs = Presentation(str(template))
    print('TEMPLATE', template.name, 'layouts', len(prs.slide_layouts))
    for idx, layout in enumerate(prs.slide_layouts):
        print('  LAYOUT', idx, repr(layout.name), 'placeholders', len(layout.placeholders))
        for ph in layout.placeholders:
            print('    PH', repr(ph.name), 'type', ph.placeholder_format.type)
    print()