from backend.components.create_ppt import create_ppt
from pptx import Presentation
slides = [
    {"title":"Test Image Slide","content":["This is a test slide with image."],"content_type":"image_text","image_index":0}
]
create_ppt(slides, 'backend', output_ppt='temp_test_output.pptx', template_name='template2.potx', images_folder='projects/Taxation/v1/slides_images')
prs = Presentation('temp_test_output.pptx')
print('slides', len(prs.slides))
for i, slide in enumerate(prs.slides, start=1):
    print('slide', i, 'shapes', len(slide.shapes))
    for shape in slide.shapes:
        ph = getattr(shape, 'placeholder_format', None)
        ph_type = getattr(ph, 'type', None)
        print('  name=', shape.name, 'shape_type=', shape.shape_type, 'ph_type=', ph_type)
        if shape.shape_type.name == 'PICTURE':
            print('    picture shape with image size', shape.image.size)
        if ph is not None:
            print('    placeholder attrs', [a for a in dir(ph) if 'insert' in a.lower() or 'picture' in a.lower()])
