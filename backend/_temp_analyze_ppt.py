from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
prs = Presentation('temp_test_output.pptx')
slide = prs.slides[0]
print('slide shapes', len(slide.shapes))
for idx, shape in enumerate(slide.shapes):
    ph = getattr(shape, 'placeholder_format', None)
    ph_type = getattr(ph, 'type', None)
    print('shape', idx, 'name=', shape.name, 'shape_type=', shape.shape_type, 'ph_type=', ph_type)
    if hasattr(shape, 'image'):
        try:
            img = shape.image
            print('  image size', img.size)
        except Exception as e:
            print('  image error', e)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print('  picture shape found')
