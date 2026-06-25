import re
import os
from pathlib import Path
import shutil
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor
 
 
# ==========================================
# COLOR PATTERNS FOR DIFFERENT TEMPLATES
# ==========================================
COLOR_PATTERNS = {
    "sample": {
        "title": RGBColor(0, 51, 102),      # Blue
        "content": RGBColor(0, 0, 0)        # Black
    },
    "sample1": {
        "title": RGBColor(0, 51, 102),      # Blue
        "content": RGBColor(0, 0, 0)        # Black
    },
    "sample2": {
        "title": RGBColor(128, 0, 0),   # custom color
        "content": RGBColor(40, 40, 40)
    },
    "sample4": {
        "title": RGBColor(255, 255, 255),   # custom color
        "content": RGBColor(255, 255, 255)
    }

}
 
# Default colors if template not found
DEFAULT_COLORS = {
    "title": RGBColor(0, 51, 102),
    "content": RGBColor(0, 0, 0)
}
 
def normalize_template_name(template_name: str):
    if not template_name:
        return None
    cleaned = str(template_name).strip().strip('"\'"')
    cleaned = Path(cleaned).stem
    return cleaned.lower()


def sanitize_template_name(template_name: str):
    if not template_name:
        return template_name
    return str(template_name).strip().strip('"\'"')


def get_colors_for_template(template_name):
    """
    Get color scheme for a specific template.
   
    Args:
        template_name: Name of the template (e.g., 'sample', 'sample3')
   
    Returns:
        Dictionary with 'title' and 'content' RGBColor values
    """
    if template_name and template_name in COLOR_PATTERNS:
        return COLOR_PATTERNS[template_name]
    return DEFAULT_COLORS
 
 
def clean_title(title: str) -> str:
    return re.sub(r'^Slide\s+\d+:\s*', '', title or "").strip()
 
 
def set_text_frame_text(text_frame, text: str, font_size: int | None = None, bold: bool | None = None, alignment=None, color=None):
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    if font_size is not None:
        p.font.size = Pt(font_size)
    if bold is not None:
        p.font.bold = bold
    if color is not None:
        p.font.color.rgb = color
        if p.runs:
            p.runs[0].font.color.rgb = color
    if alignment is not None:
        p.alignment = alignment
    return p
 
 
def estimate_paragraph_height(text: str, font_size_pt: int) -> float:
    """
    Rough height estimation (in inches) for a paragraph.
    """
    lines = max(1, len(text) // 90 + 1)
    line_height = font_size_pt * 1.25 / 72  # pt -> inches
    return lines * line_height
 
 
def split_content_to_fit(content, max_height_in, font_size_pt):
    """
    Splits content (bullets or text) into chunks that fit inside height.
    """
    chunks = []
    current = []
    current_height = 0.0
 
    for item in content:
        Para_height = estimate_paragraph_height(str(item), font_size_pt)
 
        if current_height + Para_height > max_height_in:
            chunks.append(current)
            current = [item]
            current_height = Para_height
        else:
            current.append(item)
            current_height += Para_height
 
    if current:
        chunks.append(current)
 
    return chunks
 
 
def enforce_content_character_limit(content, has_image=False):
    """
    Enforce character limits on slide content, truncating cleanly at sentence boundaries.
   
    Args:
        content: List of bullet points or content string
        has_image: Whether the slide has an image (350 char limit) or not (600 char limit)
   
    Returns:
        Content that respects character limits and ends at sentence boundaries
    """
    max_chars = 400 if has_image else 650
    max_bullets = 3 if has_image else 5
   
    if isinstance(content, list):
        total_chars = 0
        result = []
       
        # First, filter single-word bullets (already done by filter_single_word_bullets before this call)
        # Then, enforce character limit
        for item in content:
            item_str = str(item).strip()
            if not item_str:
                continue
            item_chars = len(item_str)
           
            if total_chars + item_chars <= max_chars:
                result.append(item_str)
                total_chars += item_chars
            elif total_chars < max_chars:
                # Item doesn't fit - truncate at sentence boundary
                remaining = max_chars - total_chars
                if remaining > 10:  # Only truncate if we have enough room
                    truncated = item_str[:remaining]
                   
                    # Find last sentence ending
                    last_period = truncated.rfind('.')
                    last_exclaim = truncated.rfind('!')
                    last_question = truncated.rfind('?')
                   
                    last_sentence_end = max(last_period, last_exclaim, last_question)
                   
                    if last_sentence_end != -1:
                        # Use the sentence boundary
                        truncated = truncated[:last_sentence_end + 1]
                    else:
                        # No good sentence boundary, try to end at space
                        last_space = truncated.rfind(' ')
                        if last_space != -1:
                            truncated = truncated[:last_space] + "."
                        else:
                            # Just truncate and add period if it doesn't already have one
                            if not truncated.endswith(('.', '!', '?')):
                                truncated = truncated.rstrip() + "."
                   
                    result.append(truncated)
                break # Stop adding items once limit is reached or next item doesn't fit
            else:
                break # Already at limit
       
        # Enforce maximum number of bullets (truncate extra bullets)
        if len(result) > max_bullets:
            return result[:max_bullets]
           
        return result
       
    elif isinstance(content, str):
        # For string content, truncate cleanly
        if len(content) > max_chars:
            truncated = content[:max_chars]
            # Find last sentence ending
            last_period = truncated.rfind('.')
            last_exclaim = truncated.rfind('!')
            last_question = truncated.rfind('?')
            last_sentence_end = max(last_period, last_exclaim, last_question)
           
            if last_sentence_end != -1:
                return truncated[:last_sentence_end + 1]
            else:
                last_space = truncated.rfind(' ')
                if last_space != -1:
                    return truncated[:last_space] + "."
                return truncated if truncated.endswith(('.', '!', '?')) else truncated + "."
        return content
 
 
def filter_single_word_bullets(content):
    """
    Remove single-word bullet points from slide content.
    """
    if not isinstance(content, list):
        return content
 
    filtered = []
    for item in content:
        item_str = str(item).strip()
        if not item_str:
            continue
        if len(item_str.split()) <= 1:
            continue
        filtered.append(item_str)
    return filtered
 
 
# ----------------------------
# Main function
# ----------------------------
 
def _get_best_blank_layout(prs):
    """
    Find the best blank layout from the template that preserves master slide styling.
    Most templates have multiple layouts - we want one that's blank but inherits the master.
    """
    # Try layouts in order of preference (typically blank layouts with master inheritance)
    for layout_idx in [6, 5, 7, 8, 1]:
        if layout_idx < len(prs.slide_layouts):
            return prs.slide_layouts[layout_idx]
    # Fallback to first available
    return prs.slide_layouts[0] if prs.slide_layouts else prs.slide_layouts[6]


def _find_layout_by_name(prs, keywords):
    """Return the first layout whose name contains all keywords."""
    lower_keywords = [k.lower() for k in keywords]
    for layout in prs.slide_layouts:
        layout_name = (layout.name or "").lower()
        if all(keyword in layout_name for keyword in lower_keywords):
            return layout
    return None


def _get_best_title_layout(prs):
    """Prefer a title-only layout from the template, fallback to blank."""
    return _find_layout_by_name(prs, ["title"]) or _get_best_blank_layout(prs)


def _get_best_content_layout(prs):
    """Prefer a title+content layout from the template, fallback to blank."""
    return (
        _find_layout_by_name(prs, ["title", "content"]) or
        _find_layout_by_name(prs, ["title", "body"]) or
        _find_layout_by_name(prs, ["title", "text"]) or
        _find_layout_by_placeholder_counts(prs, min_title=1, min_body=1) or
        _find_layout_by_name(prs, ["content"]) or
        _find_layout_by_name(prs, ["body"]) or
        _get_best_blank_layout(prs)
    )


def _get_best_section_header_layout(prs):
    """Prefer a section header layout from the template."""
    return (
        _find_layout_by_name(prs, ["section", "header"]) or
        _find_layout_by_name(prs, ["section"]) or
        _get_best_title_layout(prs) or
        _get_best_blank_layout(prs)
    )


def _get_best_two_content_layout(prs):
    """Prefer a two-content or comparison layout from the template."""
    exact = (
        _find_layout_by_name(prs, ["two", "content"]) or
        _find_layout_by_name(prs, ["two", "column"]) or
        _find_layout_by_name(prs, ["comparison"]) or
        _find_layout_by_placeholder_counts(prs, min_title=1, min_body=2) or
        _find_layout_by_placeholder_counts(prs, min_body=2)
    )
    if exact:
        return exact

    # Broad fallback: any layout with a title and at least 2 non-decorative placeholders
    for layout in prs.slide_layouts:
        content_count = 0
        has_title = False
        for ph in layout.placeholders:
            if not getattr(ph, "is_placeholder", False): continue
            ptype = getattr(ph.placeholder_format, "type", None)
            if ptype in (1, 3):
                has_title = True
            elif ptype not in (4, 15, 16, 17):
                content_count += 1
        if has_title and content_count >= 2:
            return layout

    return _get_best_content_layout(prs)


def _get_best_image_layout(prs):
    """Prefer a layout with title, picture and content placeholders."""
    for keywords in [
        ["title", "content", "image"],
        ["title", "content", "picture"],
        ["image", "content"],
        ["picture", "content"],
        ["title", "image"],
        ["title", "picture"],
        ["image", "text"],
        ["picture", "text"],
        ["image"],
        ["picture"]
    ]:
        match = _find_layout_by_name(prs, keywords)
        if match:
            return match

    for layout in prs.slide_layouts:
        counts = _layout_placeholder_counts(layout)
        if counts["title"] >= 1 and counts["picture"] >= 1 and counts["body"] >= 1:
            return layout

    for layout in prs.slide_layouts:
        counts = _layout_placeholder_counts(layout)
        if counts["body"] >= 1 and counts["picture"] >= 1:
            return layout

    for layout in prs.slide_layouts:
        counts = _layout_placeholder_counts(layout)
        if counts["picture"] >= 1:
            return layout

    for layout in prs.slide_layouts:
        counts = _layout_placeholder_counts(layout)
        if counts["title"] >= 1 and counts["body"] >= 2:
            return layout

    # Very broad fallback: any layout with a title and at least 2 non-decorative placeholders
    for layout in prs.slide_layouts:
        content_count = 0
        has_title = False
        for ph in layout.placeholders:
            if not getattr(ph, "is_placeholder", False): continue
            ptype = getattr(ph.placeholder_format, "type", None)
            if ptype in (1, 3): # TITLE
                has_title = True
            elif ptype not in (4, 15, 16, 17): # Ignore SUBTITLE, FOOTER, DATE, SLIDE_NUMBER
                content_count += 1
        if has_title and content_count >= 2:
            return layout

    return _get_best_content_layout(prs)


def _layout_placeholder_counts(layout):
    counts = {"title": 0, "subtitle": 0, "body": 0, "picture": 0, "object": 0}
    for ph in layout.placeholders:
        if not getattr(ph, "is_placeholder", False):
            continue
        ph_type = getattr(ph.placeholder_format, "type", None)
        name = (ph.name or "").lower()
        if ph_type == PP_PLACEHOLDER.TITLE or "title" in name:
            counts["title"] += 1
        elif ph_type == PP_PLACEHOLDER.SUBTITLE or "subtitle" in name:
            counts["subtitle"] += 1
        elif ph_type == PP_PLACEHOLDER.PICTURE or "picture" in name or "image" in name:
            counts["picture"] += 1
        elif ph_type == PP_PLACEHOLDER.BODY or "body" in name or "content" in name or "text" in name:
            counts["body"] += 1
        elif ph_type == PP_PLACEHOLDER.OBJECT:
            counts["object"] += 1
            counts["body"] += 1
    return counts


def _find_layout_by_placeholder_counts(prs, min_title=0, min_body=0, min_picture=0):
    for layout in prs.slide_layouts:
        counts = _layout_placeholder_counts(layout)
        if counts["title"] >= min_title and counts["body"] >= min_body and counts["picture"] >= min_picture:
            return layout
    return None


def _layout_has_picture_placeholder(layout):
    counts = _layout_placeholder_counts(layout)
    return counts["picture"] > 0


def _find_placeholder_by_partial_name(shapes, names):
    for shape in shapes:
        if getattr(shape, "is_placeholder", False):
            name = (shape.name or "").lower()
            if any(part in name for part in names):
                return shape
    return None


def _find_content_placeholders(shapes):
    placeholders = []
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        name = (shape.name or "").lower()
        ph_type = getattr(shape.placeholder_format, "type", None)
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            placeholders.append(shape)
        elif any(part in name for part in ["content", "body", "text", "left", "right", "column"]):
            if shape not in placeholders:
                placeholders.append(shape)
        elif ph_type not in (1, 3, 4, 15, 16, 17, 18): # Ignore Title, Subtitle, Date, Footer, SlideNum, Picture
            if shape not in placeholders:
                placeholders.append(shape)
    return placeholders


def _find_picture_placeholder(shapes):
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        name = (shape.name or "").lower()
        ph_type = getattr(shape.placeholder_format, "type", None)
        if ph_type == PP_PLACEHOLDER.PICTURE or "picture" in name or "image" in name:
            return shape
    return None


def _find_text_placeholders(shapes):
    placeholders = []
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        name = (shape.name or "").lower()
        ph_type = getattr(shape.placeholder_format, "type", None)
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            placeholders.append(shape)
        elif any(part in name for part in ["content", "body", "text", "left", "right", "column"]):
            if shape not in placeholders:
                placeholders.append(shape)
        elif ph_type not in (1, 3, 4, 15, 16, 17, 18):
            if shape not in placeholders:
                placeholders.append(shape)
    return placeholders


def _normalize_image_index(image_index):
    """Return a valid integer image index or None."""
    if image_index is None:
        return None
    if isinstance(image_index, int):
        return image_index
    if isinstance(image_index, str):
        image_index = image_index.strip()
        if image_index.isdigit():
            return int(image_index)
    try:
        return int(image_index)
    except (TypeError, ValueError):
        return None


def split_bullets_into_columns(bullets):
    """
    Split bullets into two roughly equal columns for text-only slides.
    """
    if not bullets:
        return [], []

    bullets = [str(b).strip() for b in bullets if str(b).strip()]
    if len(bullets) <= 3:
        return bullets, []

    mid = (len(bullets) + 1) // 2
    return bullets[:mid], bullets[mid:]


def _choose_layout_for_slide(prs, slide_data, content, available_images):
    ctype = slide_data.get("content_type", "content") or "content"
    image_index = _normalize_image_index(slide_data.get("image_index"))
    available_count = len(available_images) if isinstance(available_images, dict) else 0
    has_image = (
        image_index is not None
        and isinstance(available_images, dict)
        and 0 <= image_index < available_count
    )
    is_image_text = ctype in {"image_text", "image", "visual"}

    if ctype == "title":
        return _get_best_title_layout(prs)
    if ctype == "section":
        return _get_best_section_header_layout(prs)
    if is_image_text or has_image:
        return _get_best_image_layout(prs)
    if ctype == "two_column":
        return _get_best_two_content_layout(prs)
    return _get_best_content_layout(prs)


def _calculate_font_size_for_text(content, width_in, height_in, max_font_size=22, min_font_size=10):
    if not content:
        return max_font_size

    if isinstance(content, str):
        content = [content]

    text = " ".join(str(item).strip() for item in content if str(item).strip())
    if not text:
        return max_font_size

    avg_chars_per_line = max(20, int(width_in * 12))
    total_lines = sum(
        max(1, (len(str(item)) + avg_chars_per_line - 1) // avg_chars_per_line)
        for item in content
    )
    line_height_in = max_font_size * 1.12 / 72
    max_lines = max(1, int(height_in / line_height_in))

    if total_lines <= max_lines:
        return max_font_size

    computed_size = int(height_in * 72 / (total_lines * 1.12))
    return max(min_font_size, min(computed_size, max_font_size))


def _get_textbox_dimensions(shape):
    width_in = shape.width / 914400
    height_in = shape.height / 914400
    return width_in, height_in


def create_ppt(slides, audio_folder, output_ppt="output_slides.pptx", template_path=None, template_name=None, images_folder=None):
    # Allow selecting a template by full path (`template_path`) or by name
    # (`template_name`) which will be resolved inside backend/sample_ppt.
    if template_name:
        template_name = sanitize_template_name(template_name)
        root = Path(__file__).parent.parent / "sample_ppt"
        candidate = root / template_name
        if candidate.exists():
            template_path = candidate
            print(f"✅ Using selected template: {template_path}")
        else:
            # Support selection by name without extension for .pptx and .potx
            found = None
            for ext in [".pptx", ".potx"]:
                candidate = root / f"{template_name}{ext}"
                if candidate.exists():
                    found = candidate
                    break
            if found:
                template_path = found
                print(f"✅ Using selected template: {template_path}")
            else:
                print(f"❌ Template NOT found: {template_name}")
                template_path = None
    else:
        print("⚠️ No template selected. Using blank presentation.")
        template_path = None
 
    # If a template path was resolved and exists, use it, otherwise start
    # with a blank Presentation to avoid errors.
    if template_path and Path(template_path).exists():
        try:
            prs = Presentation(str(template_path))
        except Exception as exc:
            if Path(template_path).suffix.lower() == ".potx":
                with tempfile.TemporaryDirectory() as temp_dir:
                    temp_pptx = Path(temp_dir) / f"{Path(template_path).stem}.pptx"
                    shutil.copy(str(template_path), str(temp_pptx))
                    prs = Presentation(str(temp_pptx))
            else:
                raise

        print(f"✅ Template loaded with {len(prs.slide_layouts)} layouts")
        # ✅ IMPORTANT: Keep template's master slides and only remove content slides
        # This preserves the template's styling, layouts, and design
        while len(prs.slides) > 0:
            # Remove only the content (not master slide relationships)
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()
    
    # ✅ Determine which layout to use for new slides (should inherit template master design)
    blank_layout = _get_best_blank_layout(prs)
    title_layout = _get_best_title_layout(prs)
    content_layout = _get_best_content_layout(prs)
    section_layout = _get_best_section_header_layout(prs)
    two_column_layout = _get_best_two_content_layout(prs)
   
    # ✅ Collect available images if folder provided
    available_images = {}
    images_path = None
    if images_folder:
        candidate_path = Path(images_folder)
        if candidate_path.exists():
            images_path = candidate_path
        else:
            alt_path = Path(__file__).parent.parent / images_folder
            if alt_path.exists():
                images_path = alt_path

    if images_path and images_path.exists():
        for img_file in sorted(images_path.glob("*")):
            if img_file.is_file() and img_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif']:
                available_images[img_file.name] = str(img_file)
        print(f"✅ Loaded {len(available_images)} images from: {images_path}")
        if available_images:
            print(f"   Images: {list(available_images.keys())}")
    else:
        if images_folder:
            print(f"⚠️  Images folder does not exist: {images_folder}")
        else:
            print(f"⚠️  No images folder provided")
 
    # ✅ Determine whether template styling should be preserved
    normalized_template = normalize_template_name(template_name)
    template_loaded = bool(template_path and Path(template_path).exists())
    if not template_loaded:
        colors = DEFAULT_COLORS
    else:
        colors = None

    print(f"🎨 Template loaded: {template_loaded} | Template name: {normalized_template}")

 
    for idx, slide_data in enumerate(slides, start=1):
 
        title_text = clean_title(slide_data.get("title", f"Slide {idx}"))
        content = slide_data.get("content", [])
        ctype = slide_data.get("content_type", "content") or "content"
        voiceover = slide_data.get("voiceover", "")
 
        # Determine font colors based on whether a template is loaded
        title_color = colors["title"] if colors else None
        content_color = colors["content"] if colors else None
 
        # Normalize content
        if isinstance(content, str):
            content = [content]
 
        # Remove any single-word bullet points
        content = filter_single_word_bullets(content)
 
        # Determine slide layout
        chosen_layout = _choose_layout_for_slide(prs, slide_data, content, available_images) or blank_layout
        # Check if this is a title-only slide (cover slide)
        # First slide is ALWAYS title-only (centered)
        is_title_only = (idx == 1) or ctype == "title" or (not content or all(not str(item).strip() for item in content))
 
        if is_title_only:
            # ----------------------------
            # TITLE-ONLY SLIDE (CENTERED)
            # ----------------------------
            slide = prs.slides.add_slide(chosen_layout)
            title_shape = slide.shapes.title if slide.shapes.title else None

            if title_shape:
                set_text_frame_text(
                    title_shape.text_frame,
                    title_text,
                    font_size=54,
                    bold=True,
                    alignment=PP_ALIGN.CENTER,
                    color=title_color
                )
            else:
                slide_height = prs.slide_height
                slide_width = prs.slide_width
 
                title_box = slide.shapes.add_textbox(
                    Inches(0.8),
                    (slide_height - Inches(2)) / 2,  # Vertical center
                    slide_width - Inches(1.6),
                    Inches(2)
                )
 
                tf = title_box.text_frame
                set_text_frame_text(
                    tf,
                    title_text,
                    font_size=54,
                    bold=True,
                    alignment=PP_ALIGN.CENTER,
                    color=title_color
                )
                tf.auto_size = True
 
            # ----------------------------
            # NOTES
            # ----------------------------
            slide.notes_slide.notes_text_frame.text = (
                f"{title_text}\n\n"
                f"{voiceover}"
            )
 
        else:
            # ----------------------------
            # REGULAR SLIDES WITH CONTENT
            # ----------------------------
            # Content zone measurements
            content_top = Inches(2.8)
            content_height = prs.slide_height - Inches(3.5)
 
            content_font_size = 22
 
            image_index = _normalize_image_index(slide_data.get("image_index"))
            available_count = len(available_images) if isinstance(available_images, dict) else 0
            has_image = (
                image_index is not None
                and isinstance(available_images, dict)
                and 0 <= image_index < available_count
            )

            # Log image processing
            if image_index is not None:
                print(f"  Slide {idx}: image_index={image_index}, available_images={available_count}, has_image={has_image}")
                if not has_image and available_count > 0:
                    if not isinstance(image_index, int):
                        print(f"    → image_index is not an integer-like value (actual type: {type(slide_data.get('image_index'))})")
                    elif image_index < 0 or image_index >= available_count:
                        print(f"    → image_index out of range [0, {available_count-1}]")
 
            slide = prs.slides.add_slide(chosen_layout)
            title_shape = slide.shapes.title if slide.shapes.title else None

            # ----------------------------
            # TITLE BOX (AUTO SAFE)
            # ----------------------------
            if title_shape:
                set_text_frame_text(
                    title_shape.text_frame,
                    title_text,
                    font_size=36,
                    bold=True,
                    alignment=PP_ALIGN.LEFT,
                    color=title_color
                )
            else:
                title_box = slide.shapes.add_textbox(
                    Inches(0.8),
                    Inches(1.2),
                    prs.slide_width - Inches(1.6),
                    Inches(1.4)  # taller title box
                )
 
                tf = title_box.text_frame
                set_text_frame_text(
                    tf,
                    title_text,
                    font_size=36,
                    bold=True,
                    alignment=PP_ALIGN.LEFT,
                    color=title_color
                )
                tf.auto_size = True
 
            # ----------------------------
            # CONTENT BOX
            # ----------------------------
            enforced_chunk = enforce_content_character_limit(content, has_image=has_image)
 
            def add_bullet_list_to_frame(frame, bullets, shape=None):
                if isinstance(bullets, str):
                    bullets = [bullets]
                elif bullets is None:
                    bullets = []
                frame.clear()
                frame.word_wrap = True
                width_in, height_in = _get_textbox_dimensions(shape) if shape is not None else (prs.slide_width / 914400 - 1.6, content_height / 914400)
                font_size = _calculate_font_size_for_text(bullets, width_in, height_in, max_font_size=22, min_font_size=10)
                for i, bullet in enumerate(bullets):
                    para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                    para.text = str(bullet)
                    if font_size:
                        para.font.size = Pt(font_size)
                    if content_color is not None:
                        para.font.color.rgb = content_color
                        if para.runs:
                            para.runs[0].font.color.rgb = content_color
                    para.level = 0
 
            content_placeholders = _find_text_placeholders(slide.shapes)
            picture_placeholder = _find_picture_placeholder(slide.shapes)
            image_filenames = list(available_images.keys()) if available_images else []
            needs_image_layout = slide_data.get("content_type") in {"image_text", "image", "visual"} or has_image

            if needs_image_layout:
                image_filename = image_filenames[image_index] if image_index is not None and 0 <= image_index < len(image_filenames) else None
                if image_filename:
                    image_path = available_images[image_filename]
                    try:
                        if picture_placeholder is not None:
                            if hasattr(picture_placeholder, "insert_picture"):
                                picture_placeholder.insert_picture(image_path)
                            else:
                                slide.shapes.add_picture(
                                    image_path,
                                    picture_placeholder.left,
                                    picture_placeholder.top,
                                    width=picture_placeholder.width,
                                    height=picture_placeholder.height
                                )
                        elif len(content_placeholders) >= 2:
                            image_target = content_placeholders[1]
                            slide.shapes.add_picture(
                                image_path,
                                image_target.left,
                                image_target.top,
                                width=image_target.width,
                                height=image_target.height
                            )
                        else:
                            slide.shapes.add_picture(
                                image_path,
                                Inches(0.8),
                                content_top,
                                width=prs.slide_width - Inches(1.6),
                                height=content_height
                            )
                        print(f"  ✓ Slide {idx}: Added image #{image_index} ({image_filename})")
                    except Exception as e:
                        print(f"  ⚠️  Slide {idx}: Failed to add image #{image_index} ({image_filename}) - {e}")
                elif picture_placeholder is not None:
                    # No file, leave image placeholder blank for manual replacement
                    print(f"  ⚠️  Slide {idx}: No image file available, preserving image placeholder")
                else:
                    print(f"  ⚠️  Slide {idx}: No picture placeholder found for image slide layout")

                if content_placeholders:
                    if len(content_placeholders) >= 2 and slide_data.get("content_type") == "two_column":
                        left_bullets, right_bullets = split_bullets_into_columns(enforced_chunk)
                        add_bullet_list_to_frame(content_placeholders[0].text_frame, left_bullets, shape=content_placeholders[0])
                        add_bullet_list_to_frame(content_placeholders[1].text_frame, right_bullets, shape=content_placeholders[1])
                    else:
                        target_frame = content_placeholders[0].text_frame
                        add_bullet_list_to_frame(target_frame, enforced_chunk, shape=content_placeholders[0])
                elif picture_placeholder is not None:
                    # Add a text box only if image placeholder exists but no text placeholder
                    content_box = slide.shapes.add_textbox(
                        picture_placeholder.left + picture_placeholder.width + Inches(0.2),
                        content_top,
                        prs.slide_width - (picture_placeholder.left + picture_placeholder.width + Inches(1.0)),
                        content_height
                    )
                    add_bullet_list_to_frame(content_box.text_frame, enforced_chunk, shape=content_box)
                else:
                    content_box = slide.shapes.add_textbox(
                        Inches(0.8),
                        content_top,
                        prs.slide_width / 2 - Inches(1.0),
                        content_height
                    )
                    add_bullet_list_to_frame(content_box.text_frame, enforced_chunk, shape=content_box)
            elif content_placeholders:
                if len(content_placeholders) >= 2 and slide_data.get("content_type") == "two_column":
                    left_bullets, right_bullets = split_bullets_into_columns(enforced_chunk)
                    add_bullet_list_to_frame(content_placeholders[0].text_frame, left_bullets, shape=content_placeholders[0])
                    add_bullet_list_to_frame(content_placeholders[1].text_frame, right_bullets, shape=content_placeholders[1])
                else:
                    target_frame = content_placeholders[0].text_frame
                    add_bullet_list_to_frame(target_frame, enforced_chunk, shape=content_placeholders[0])
            else:
                use_two_columns = (
                    isinstance(enforced_chunk, list)
                    and len(enforced_chunk) > 3
                    and ctype == "two_column"
                )
                if use_two_columns:
                    left_bullets, right_bullets = split_bullets_into_columns(enforced_chunk)
                    left_box = slide.shapes.add_textbox(
                        Inches(0.8),
                        content_top,
                        prs.slide_width / 2 - Inches(1.1),
                        content_height
                    )
                    right_box = slide.shapes.add_textbox(
                        prs.slide_width / 2 + Inches(0.1),
                        content_top,
                        prs.slide_width / 2 - Inches(1.1),
                        content_height
                    )
                    add_bullet_list_to_frame(left_box.text_frame, left_bullets, shape=left_box)
                    add_bullet_list_to_frame(right_box.text_frame, right_bullets, shape=right_box)
                else:
                    content_box = slide.shapes.add_textbox(
                        Inches(0.8),
                        content_top,
                        prs.slide_width - Inches(1.6),
                        content_height
                    )
                    tf = content_box.text_frame
                    add_bullet_list_to_frame(tf, enforced_chunk, shape=content_box)
 
            # ----------------------------
            # NOTES
            # ----------------------------
            slide.notes_slide.notes_text_frame.text = (
                f"{title_text}\n\n"
                f"{voiceover}\n\n"
                f"Audio: {os.path.join(audio_folder, f'slide_{idx}.mp3')}"
            )
 
    
    if os.path.exists(output_ppt):
        os.remove(output_ppt)

    prs.save(output_ppt)
    final_slide_count = len(prs.slides)
 
    print(f"✅ PPT created with {final_slide_count} slides")
    return final_slide_count